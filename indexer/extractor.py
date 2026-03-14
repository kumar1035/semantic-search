# indexer/extractor.py

import os
import json
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook


class Extractor:
    """
    Extracts raw text content from different file types.
    Each file type has its own extraction method.
    """

    def extract(self, filepath):
        """
        Main dispatcher — picks the right extraction method based on file extension.
        """
        handlers = {
            ".pdf": self.extract_pdf,
            ".docx": self.extract_docx,
            ".pptx": self.extract_pptx,
            ".xlsx": self.extract_xlsx,
            ".ipynb": self.extract_ipynb,
            ".txt": self.extract_text,
            ".md": self.extract_text,
            ".py": self.extract_text,
            ".js": self.extract_text,
        }

        try:
            ext = os.path.splitext(filepath)[1].lower()
            handler = handlers.get(ext)
            if handler:
                return handler(filepath)
            else:
                print(f"Warning: Unrecognized file extension: {ext}")
                return ""
        except Exception as e:
            print(f"Error extracting text from {filepath}: {e}")
            return ""

    def extract_pdf(self, filepath):
        """Extract text from a PDF file using PyMuPDF."""
        doc = fitz.open(filepath)
        pages = []
        for page in doc:
            pages.append(page.get_text())
        doc.close()
        return "\n".join(pages)

    def extract_docx(self, filepath):
        """Extract text from a Word document using python-docx."""
        doc = Document(filepath)
        paragraphs = []
        for para in doc.paragraphs:
            paragraphs.append(para.text)
        return "\n".join(paragraphs)

    def extract_pptx(self, filepath):
        """Extract text from a PowerPoint file using python-pptx."""
        prs = Presentation(filepath)
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        lines.append(para.text)
        return "\n".join(lines)

    def extract_xlsx(self, filepath):
        """Extract text from an Excel file using openpyxl."""
        wb = load_workbook(filepath, data_only=True)
        rows = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            for row in sheet.iter_rows():
                cells = []
                for cell in row:
                    if cell.value is not None:
                        cells.append(str(cell.value))
                rows.append(" ".join(cells))
        return "\n".join(rows)

    def extract_ipynb(self, filepath):
        """Extract text from a Jupyter notebook (.ipynb) file."""
        with open(filepath, "r", encoding="utf-8") as f:
            notebook = json.load(f)
        cells = []
        for cell in notebook["cells"]:
            cell_text = "".join(cell["source"])
            cells.append(cell_text)
        return "\n".join(cells)

    def extract_text(self, filepath):
        """Extract text from plain text files (.txt, .md, .py, .js, etc.)"""
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()


# --- Test it ---
if __name__ == "__main__":
    import sys

    extractor = Extractor()

    if len(sys.argv) > 1:
        filepath = sys.argv[1]
        text = extractor.extract(filepath)
        print(f"Extracted {len(text)} characters from {filepath}")
        print(f"Preview:\n{text[:500]}")
    else:
        print("Usage: python -m indexer.extractor <filepath>")